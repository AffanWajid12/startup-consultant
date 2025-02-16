import streamlit as st
import requests
import time
import json
import os
import plotly.graph_objects as go
import plotly.express as px
import pandas as pd
from dataclasses import dataclass
from typing import Dict, List, Optional, Any
from enum import Enum
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

# ----- LangChain & CrewAI Imports -----
from langchain_community.document_loaders import TextLoader
from pytrends.request import TrendReq
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain.schema import Document
from openai import OpenAI
from crewai import Agent, Task, Crew

# ----- PDF Generation Imports -----
from reportlab.lib.pagesizes import letter, inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.pdfgen import canvas

# ----- Additional LangChain Imports -----
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings

# Import LangChain's LLM base class from the correct location (Pydantic-based)
from langchain.llms.base import LLM

# ----- Load Environment Variables -----
load_dotenv()

# =============================================================================
# Module 1: Startup Consultant (Data Fetching + RAG Pipeline)
# =============================================================================

@st.cache_data
def fetch_google_trends_data(keyword="startups"):
    pytrends = TrendReq(retries=3, backoff_factor=0.1)
    time.sleep(2)  # Delay before building payload
    pytrends.build_payload([keyword], timeframe='now 7-d', geo='US')
    time.sleep(2)  # Delay before fetching data
    trends = pytrends.interest_over_time()
    return trends.to_json()

@st.cache_data
def fetch_newsapi_data(query="startups"):
    # Fetch the NewsAPI key from the .env file
    newsapi_key = os.getenv("NEWSAPI_KEY")
    url = f"https://newsapi.org/v2/everything?q={query}&apiKey={newsapi_key}"
    response = requests.get(url)
    return response.json().get("articles", [])

@st.cache_data
def load_combined_data():
    google_data = fetch_google_trends_data()
    news_data = fetch_newsapi_data()
    combined_docs = [
        Document(page_content=f"Google Trends Data: {google_data}"),
        Document(page_content=f"News Articles: {news_data}")
    ]
    return combined_docs

def create_vector_store(docs):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_documents(docs, embeddings)

# Prepare documents and build retriever
data_docs = load_combined_data()
doc_chunks = RecursiveCharacterTextSplitter(chunk_size=500).split_documents(data_docs)
vectorstore = create_vector_store(doc_chunks)
retriever = vectorstore.as_retriever(search_kwargs={"k": 10})

# =============================================================================
# Set up AIML API client (all calls go through https://api.aimlapi.com/v1)
# =============================================================================

# Fetch the AIML API key from the .env file
aiml_api_key = os.getenv("AIML_API_KEY")
base_url = "https://api.aimlapi.com/v1"
model_name = "deepseek-ai/deepseek-llm-67b-chat"
api = OpenAI(api_key=aiml_api_key, base_url=base_url)

# -----------------------------------------------------------------------------
# New AIMLLLMWrapper class that inherits from LangChain's LLM (as a Pydantic model)
# -----------------------------------------------------------------------------
class AIMLLLMWrapper(LLM):
    api: Any
    model: str
    temperature: float = 0.5
    max_tokens: int = 2000

    @property
    def _llm_type(self) -> str:
        return "aiml"

    def _call(self, prompt: str, stop: Optional[List[str]] = None) -> str:
        response = self.api.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=self.temperature,
            max_tokens=self.max_tokens,
        )
        return response.choices[0].message.content

    async def _acall(self, prompt: str, stop: Optional[List[str]] = None) -> str:
        raise NotImplementedError("Async call not supported.")

# Instantiate our LLM wrapper using the API key from .env
llm = AIMLLLMWrapper(api=api, model=model_name, temperature=0.3)

# Define prompt for Startup Consultant
system_prompt = (
    "You are an AI Startup Consultant specializing in market research, idea validation, and business insights. "
    "Use the provided context to answer user queries related to startup ideas, trends, competitors, and growth potential. "
    "If unsure, respond with 'I don't know' instead of guessing. "
    "\n\n{context}"
)
prompt = ChatPromptTemplate.from_messages([
    ("system", system_prompt),
    ("human", "{input}")
])

def startup_consultant_tab():
    st.header("Module 1: Startup Consultant")
    query = st.text_input("🗨️ Enter your startup-related query:")
    if st.button("Submit Query", key="consultant_submit"):
        with st.spinner("Fetching insights..."):
            question_chain = create_stuff_documents_chain(llm, prompt)
            rag_chain = create_retrieval_chain(retriever, question_chain)
            response = rag_chain.invoke({"input": query})
            result = response.get("answer", "No answer returned.")
            st.session_state.model_1_output = result  # Save for Module 2
            st.write("**Response from Startup Consultant:**")
            st.write(result)

# =============================================================================
# Module 2: Business Model Generation
# =============================================================================

def generate_business_model(scenario):
    system_prompt_biz = (
        "You are a top-tier financial strategist and business planner. Based on the given scenario, "
        "create a comprehensive analysis that includes: "
        "\n\n1. REVENUE MODEL: Recommend the most suitable model (subscription, marketplace, SaaS, etc.) with detailed pricing structure. "
        "\n2. FINANCIAL FORECAST: Provide 3-year projections of revenue, costs, and profitability with quarterly breakdown. "
        "\n3. BREAK-EVEN ANALYSIS: Calculate and explain when the business will become profitable. "
        "\n4. MONETIZATION STRATEGIES: Suggest three distinct approaches with their respective advantages and disadvantages. "
        "\n5. ONE-PAGE BUSINESS PLAN: Create a concise, investor-ready summary that highlights the value proposition, market opportunity, "
        "revenue model, and path to profitability. "
    )
    completion = api.chat.completions.create(
        model=model_name,
        messages=[
            {"role": "system", "content": system_prompt_biz},
            {"role": "user", "content": scenario},
        ],
        temperature=0.5,
        max_tokens=2000,
    )
    return completion.choices[0].message.content

def business_model_tab():
    st.header("Module 2: Business Model Generation")
    if st.button("Generate Business Model", key="biz_model_submit"):
        if not st.session_state.get("model_1_output"):
            st.error("Please run the Startup Consultant module first.")
        else:
            scenario = st.session_state.model_1_output
            with st.spinner("Generating business model..."):
                biz_response = generate_business_model(scenario)
                st.session_state.business_model_output = biz_response  # Save for Module 3
                st.write("**Business Model Analysis:**")
                st.write(biz_response)

# =============================================================================
# Module 3: Pitch Deck Generation
# =============================================================================

# --- Define Data Structures and Classes for Pitch Deck Generation ---

class VisualType(Enum):
    CHART = "chart"
    IMAGE = "image"
    TIMELINE = "timeline"
    GRAPH = "graph"

@dataclass
class SlideContent:
    title: str
    content: str
    visual_type: Optional[VisualType] = None
    visual_data: Optional[Dict[str, Any]] = None

@dataclass
class PitchDeckData:
    company_name: str
    problem_statement: str
    solution: str
    market_size: float
    revenue_model: Dict[str, float]
    roadmap: List[Dict[str, Any]]
    team: List[Dict[str, str]]
    traction: str
    future_outlook: str

class AIContentGenerator:
    def generate_elevator_pitch(self, pitch_data: PitchDeckData) -> str:
        prompt_text = f"""Create a compelling elevator pitch for {pitch_data.company_name}.
Problem: {pitch_data.problem_statement}
Solution: {pitch_data.solution}
Market Size: ${pitch_data.market_size:,.2f}
"""
        try:
            response = api.chat.completions.create(
                model=model_name,
                messages=[{"role": "user", "content": prompt_text}],
                temperature=0.5,
                max_tokens=500,
            )
            return response.choices[0].message.content
        except Exception as e:
            st.error(f"Error generating elevator pitch: {e}")
            return "Error generating elevator pitch. Please check the logs."

    def generate_executive_summary(self, pitch_data: PitchDeckData) -> str:
        team_description = "Experienced team"
        if pitch_data.team:
            try:
                roles = [member.get('role', 'undefined role') for member in pitch_data.team]
                team_description = f"{len(pitch_data.team)} members with {roles[0]} leadership"
            except (KeyError, IndexError):
                team_description = "Core team in place"

        prompt_text = f"""Create a concise executive summary for {pitch_data.company_name} covering:
- Problem: {pitch_data.problem_statement}
- Solution: {pitch_data.solution} 
- Market Size: ${pitch_data.market_size:,.2f}
- Revenue Model: {json.dumps(pitch_data.revenue_model)}
- Key Traction: {pitch_data.traction}
- Roadmap Highlights: {pitch_data.roadmap[:2]}
- Team Strength: {team_description}
"""
        try:
            response = api.chat.completions.create(
                model=model_name,
                messages=[{"role": "user", "content": prompt_text}],
                temperature=0.3,
                max_tokens=1000,
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"Error generating executive summary: {e}"

class VisualizationGenerator:
    def create_market_size_chart(self, market_data):
        fig = go.Figure()
        fig.add_trace(go.Bar(
            x=list(market_data.keys()),
            y=list(market_data.values()),
            text=[f'${x / 1e6:.1f}M' for x in market_data.values()],
            marker_color=['#2962FF', '#546E7A', '#00C853']
        ))
        fig.update_layout(
            template='plotly_white',
            title_font_size=20,
            hoverlabel=dict(bgcolor="white")
        )
        return fig

    def create_revenue_projection(self, revenue_data: Dict[str, float]) -> go.Figure:
        df = pd.DataFrame(list(revenue_data.items()), columns=['Stream', 'Revenue'])
        fig = px.pie(df, values='Revenue', names='Stream',
                     title='Revenue Distribution',
                     color_discrete_sequence=['#2962FF', '#00C853'])
        return fig

    def create_roadmap_timeline(self, milestones: List[Dict[str, Any]]) -> go.Figure:
        df = pd.DataFrame(milestones)
        fig = px.timeline(df, x_start='start_date', x_end='end_date',
                          y='milestone', title='Product Roadmap',
                          color_discrete_sequence=['#2962FF'])
        return fig

    def create_solution_diagram(self):
        fig = go.Figure(go.Indicator(
            mode="number+gauge",
            value=50,
            domain={'x': [0.1, 1], 'y': [0.1, 1]},
            title={'text': "Cost Savings %"},
            gauge={'shape': "bullet"}
        ))
        return fig

    def create_financial_forecast(self):
        years = [2024, 2025, 2026]
        revenue = [1.2, 3.5, 8.0]  # In millions
        fig = px.line(x=years, y=revenue, title="Revenue Projection (Millions USD)")
        fig.update_traces(line_color='#2962FF')
        return fig


class PresentationBuilder:
    def __init__(self, template_path: str = None):
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self._apply_theme()
        self.pdf_story = []

    def _apply_theme(self):
        for slide in self.prs.slides:
            self._set_slide_background(slide)

    def _set_slide_background(self, slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)

    def add_title_slide(self, company_name: str, subtitle: str):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)
        title = slide.shapes.title
        title.text = company_name
        title.text_frame.paragraphs[0].font.size = Pt(60)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 70, 122)
        title.text_frame.alignment = PP_ALIGN.CENTER
        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = subtitle
        subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
        subtitle_placeholder.text_frame.paragraphs[0].font.italic = True
        subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        subtitle_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(name='TitleStyle', parent=styles['h1'], alignment=TA_CENTER)
        subtitle_style = ParagraphStyle(name='SubtitleStyle', parent=styles['h3'], alignment=TA_CENTER)
        self.pdf_story.append(Paragraph(company_name, title_style))
        self.pdf_story.append(Paragraph(subtitle, subtitle_style))
        self.pdf_story.append(Spacer(1, 0.5 * inch))

    def add_content_slide(self, content: SlideContent):
        blank_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[-1]
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_slide_background(slide)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
        tf_title = title_box.text_frame
        tf_title.text = content.title.upper()
        p_title = tf_title.paragraphs[0]
        p_title.font.bold = True
        p_title.font.size = Pt(44)
        p_title.font.color.rgb = RGBColor(0, 70, 122)
        p_title.alignment = PP_ALIGN.CENTER

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        for line in content.content.split('\n'):
            p = tf_content.add_paragraph()
            p.text = line
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(10)

        if content.visual_type and content.visual_data:
            visual_left = Inches(9)
            visual_top = Inches(1.5)
            visual_width = Inches(6)
            if content.visual_type == VisualType.CHART:
                self._add_chart(slide, content.visual_data, visual_left, visual_top, visual_width)
            elif content.visual_type == VisualType.TIMELINE:
                self._add_timeline(slide, content.visual_data, visual_left, visual_top, visual_width)
            elif content.visual_type == VisualType.IMAGE:
                self._add_image(slide, content.visual_data, visual_left, visual_top, visual_width)

        styles = getSampleStyleSheet()
        heading_style = styles['h2']
        body_style = styles['Normal']
        heading_style.alignment = TA_LEFT  # Align headings to left in PDF
        self.pdf_story.append(Paragraph(content.title, heading_style))
        for line in content.content.split('\n'):
            self.pdf_story.append(Paragraph(line, body_style))
        self.pdf_story.append(Spacer(1, 0.2 * inch))  # Add a smaller space after each content slide section in PDF
        self.pdf_story.append(PageBreak())  # Add page break after each content slide for PDF

    def _add_chart(self, slide, visual_data, visual_left, visual_top, visual_width):
        if 'chart' in visual_data and visual_data['chart'] is not None:
            chart_filepath = "temp_chart.png"
            chart = visual_data['chart']
            chart.update_layout(
                width=800,
                height=600,
                margin=dict(l=20, r=20, t=40, b=20),
                title_font_size=24,
                plot_bgcolor='rgba(0,0,0,0)',
                paper_bgcolor='rgba(0,0,0,0)'
            )
            chart.write_image(chart_filepath)
            try:
                slide.shapes.add_picture(chart_filepath, visual_left, visual_top, width=visual_width)
            except Exception as e:
                st.error(f"Error adding chart to slide: {e}")
            finally:
                if os.path.exists(chart_filepath):
                    os.remove(chart_filepath)

    def _add_timeline(self, slide, visual_data, visual_left, visual_top, visual_width):
        if 'chart' in visual_data and visual_data['chart'] is not None:
            timeline_filepath = "temp_timeline.png"
            timeline_chart = visual_data['chart']
            timeline_chart.write_image(timeline_filepath)
            try:
                slide.shapes.add_picture(timeline_filepath, visual_left, visual_top, width=visual_width)
            except Exception as e:
                st.error(f"Error adding timeline to slide: {e}")
            finally:
                if os.path.exists(timeline_filepath):
                    os.remove(timeline_filepath)

    def _add_image(self, slide, visual_data, left, top, width):
        try:
            slide.shapes.add_picture(visual_data['image'], left, top, width=width)
        except Exception as e:
            st.error(f"Error adding image: {e}")

    def save_pptx(self, filename: str):
        self.prs.save(filename)

    def save_pdf(self, filename: str, company_name: str):
        doc = SimpleDocTemplate(filename, pagesize=letter)
        doc.build(self.pdf_story,
                  onFirstPage=lambda canvas, doc: self._add_pdf_header_footer(canvas, doc, company_name),
                  onLaterPages=lambda canvas, doc: self._add_pdf_header_footer(canvas, doc, company_name))

    def _add_pdf_header_footer(self, canvas, doc, company_name):
        canvas.saveState()
        canvas.setFont('Times-Roman', 9)
        canvas.drawString(inch, 0.75 * inch, f"Page {doc.page}")
        canvas.drawString(letter[0] - inch, 0.75 * inch, company_name)  # Company name as footer right
        canvas.restoreState()


class PitchDeckGenerator:
    def __init__(self):
        self.content_generator = AIContentGenerator()
        self.viz_generator = VisualizationGenerator()
        self.presentation_builder = PresentationBuilder()

    def generate_pitch_deck(self, pitch_data: PitchDeckData) -> tuple[str, str]:
        # Generate elevator pitch and executive summary using ChatOllama
        elevator_pitch = self.content_generator.generate_elevator_pitch(pitch_data)
        exec_summary = self.content_generator.generate_executive_summary(pitch_data)

        # Create presentation title slide with elevator pitch as subtitle
        self.presentation_builder.add_title_slide(
            pitch_data.company_name, elevator_pitch)

        # Generate additional slides
        slides = self._generate_slides(pitch_data, exec_summary)
        for slide in slides:
            self.presentation_builder.add_content_slide(slide)

        pptx_output_path = f"pitch_deck_{pitch_data.company_name.lower().replace(' ', '_')}.pptx"
        pdf_output_path = f"pitch_deck_{pitch_data.company_name.lower().replace(' ', '_')}.pdf"

        self.presentation_builder.save_pptx(pptx_output_path)
        self.presentation_builder.save_pdf(pdf_output_path, pitch_data.company_name)  # Save PDF as well

        return pptx_output_path, pdf_output_path  # Return both paths

    def _generate_slides(self, pitch_data: PitchDeckData, exec_summary: str) -> List[SlideContent]:
        slides = []

        if "[**Specify" in exec_summary:
            raise ValueError("AI failed to generate proper executive summary")
        if len(pitch_data.solution) < 50:
            raise ValueError("Solution description too vague")

        # Problem slide
        slides.append(SlideContent(
            title="The Problem",
            content=f"**Key Pain Points:**\n{pitch_data.problem_statement}",
            # visual_type=VisualType.IMAGE,
            # visual_data={"image": "path/to/problem_icon.png"}  # Adjust image path
        ))

        # Solution slide
        slides.append(SlideContent(
            title="Our Solution",
            content=f"**{pitch_data.company_name}'s Innovation:**\n{pitch_data.solution}\n\n**Key Benefits:**\n- 50% cost savings vs competitors\n- 98% customer satisfaction",
            visual_type=VisualType.CHART,
            visual_data={"chart": self.viz_generator.create_solution_diagram()}
        ))

        # Market size slide
        market_chart = self.viz_generator.create_market_size_chart({
            'TAM': pitch_data.market_size,
            'SAM': pitch_data.market_size * 0.6,
            'SOM': pitch_data.market_size * 0.3
        })
        slides.append(SlideContent(
            title="Market Opportunity",
            content=f"**${pitch_data.market_size / 1e9:.1f}B Total Addressable Market**\n22% CAGR projected (2024-2029)",
            visual_type=VisualType.CHART,
            visual_data={"chart": market_chart}
        ))

        # Revenue projection slide
        rev_chart = self.viz_generator.create_revenue_projection(pitch_data.revenue_model)
        slides.append(SlideContent(
            title="Business Model",
            content="**Revenue Streams:**\n" + "\n".join(
                [f"- {k}: {v}%" for k, v in pitch_data.revenue_model.items()]),
            visual_type=VisualType.CHART,
            visual_data={"chart": rev_chart}
        ))

        # Roadmap timeline slide
        roadmap_chart = self.viz_generator.create_roadmap_timeline(pitch_data.roadmap)
        slides.append(SlideContent(
            title="Product Roadmap",
            content="Key Milestones & Timeline",
            visual_type=VisualType.TIMELINE,
            visual_data={"chart": roadmap_chart}
        ))

        # Traction slide
        slides.append(SlideContent(
            title="Traction & Validation",
            content=f"**Early Success:**\n{pitch_data.traction}\n\n**Key Metrics:**\n- 80% Pilot Retention\n- 4.9/5 Customer Rating"
        ))

        # Financials slide
        slides.append(SlideContent(
            title="Financial Projections",
            content="3-Year Growth Outlook",
            visual_type=VisualType.CHART,
            visual_data={"chart": self.viz_generator.create_financial_forecast()}
        ))

        # Team slide
        slides.append(SlideContent(
            title="Leadership Team",
            content="\n".join([f"- {m['name']} ({m['role']})" for m in pitch_data.team]),
            # visual_type=VisualType.IMAGE,
            # visual_data={"image": "path/to/team_photo.png"}  # Adjust image path
        ))

        # Executive summary slide
        slides.append(SlideContent(
            title="Executive Summary",
            content=exec_summary
        ))

        # Conclusion slide
        slides.append(SlideContent(
            title="Next Steps",
            content="**Investment Ask:** $2M Seed Round\n**Key Milestones:**\n- Expand to 3 new cities\n- Launch mobile app"
        ))
        return slides


def pitch_deck_tab():
    st.header("Module 3: Pitch Deck Generation")
    st.write(
        "This section uses data from previous modules to prefill pitch deck details. You may adjust the inputs as needed.")

    # Pre-populate fields from previous modules if available
    default_company_name = st.session_state.get("company_name", "Your Company")
    default_problem = st.session_state.get("model_1_output", "Describe the key pain points here.")
    default_solution = st.session_state.get("business_model_output",
                                            "Describe your solution and business model here. Include key features, revenue approach, etc.")
    default_market_size = 1_000_000_000.0  # Example default (1B USD)
    default_revenue_streams = '{"Subscription": 70, "Licensing": 30}'
    default_roadmap = '[{"milestone": "MVP", "start_date": "2024-03", "end_date": "2024-06"}]'
    default_team = '[{"name": "John Doe", "role": "CEO"}]'
    default_traction = st.session_state.get("business_model_output", "Mention early customer/adoption metrics here.")
    default_future = "Outline future growth opportunities here."

    with st.form("pitch_deck_form"):
        company_name = st.text_input("Company Name", value=default_company_name)
        problem = st.text_area("Problem Statement", value=default_problem)
        solution = st.text_area("Solution", value=default_solution)
        market_size = st.number_input("Total Addressable Market Size (USD)", min_value=0.0, value=default_market_size)
        st.subheader("Revenue Model")
        revenue_streams = st.text_area(
            "Revenue Streams (Enter as JSON)",
            value=default_revenue_streams
        )
        st.subheader("Roadmap")
        roadmap = st.text_area(
            "Roadmap Milestones (Enter as JSON)",
            value=default_roadmap
        )
        st.subheader("Team")
        team = st.text_area(
            "Team Members (Enter as JSON)",
            value=default_team
        )
        traction = st.text_area("Traction Metrics", value=default_traction)
        future_outlook = st.text_area("Future Outlook", value=default_future)

        submitted = st.form_submit_button("Generate Pitch Deck")

        if submitted:
            try:
                st.write("**Debugging Inputs:**")
                st.write("Revenue Streams:", revenue_streams)
                st.write("Roadmap:", roadmap)
                st.write("Team:", team)
                st.write("--- End of Debugging ---")
                pitch_data = PitchDeckData(
                    company_name=company_name,
                    problem_statement=problem,
                    solution=solution,
                    market_size=float(market_size),
                    revenue_model=json.loads(revenue_streams),
                    roadmap=json.loads(roadmap),
                    team=json.loads(team),
                    traction=traction,
                    future_outlook=future_outlook
                )
                generator = PitchDeckGenerator()
                pptx_path, pdf_path = generator.generate_pitch_deck(pitch_data)  # Get both paths
                st.session_state.output_path_pptx = pptx_path  # Store pptx path
                st.session_state.output_path_pdf = pdf_path  # Store pdf path
                st.success("Pitch deck (PPTX & PDF) generated successfully!")
            except Exception as e:
                st.error(f"Error generating pitch deck: {str(e)}")

    if st.session_state.get("output_path_pptx"):
        with open(st.session_state.output_path_pptx, "rb") as file:
            st.download_button(
                label="Download Pitch Deck (PPTX)",
                data=file,
                file_name=st.session_state.output_path_pptx,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

    if st.session_state.get("output_path_pdf"):
        with open(st.session_state.output_path_pdf, "rb") as file:
            st.download_button(
                label="Download Pitch Deck (PDF)",
                data=file,
                file_name=st.session_state.output_path_pdf,
                mime="application/pdf"
            )


# =============================================================================
# Main App: Integrate All Modules Using Tabs
# =============================================================================

def main():
    st.title("Integrated AI Startup Builder & Pitch Generator")
    st.write(
        "Follow the tabs below to progress through: Startup Consultation → Business Model Generation → Pitch Deck Creation")

    tabs = st.tabs(["Startup Consultant", "Business Model", "Pitch Deck"])
    with tabs[0]:
        startup_consultant_tab()
    with tabs[1]:
        business_model_tab()
    with tabs[2]:
        pitch_deck_tab()


if __name__ == "__main__":
    main()
